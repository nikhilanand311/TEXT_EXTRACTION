import streamlit as st
import pytesseract
from PIL import Image, ImageDraw, ImageFont
import fitz  # PyMuPDF for PDF text extraction
import io
import docx
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import os
import re
from dotenv import load_dotenv
from transformers import pipeline
import spacy

# Load environment variables from .env file
load_dotenv()

# Initialize spaCy model
nlp = spacy.load("en_core_web_sm")

# Initialize Hugging Face summarization pipeline
summarizer = pipeline("summarization")

# Retrieve the API key from environment variable
api_key = os.getenv('OPENAI_API_KEY')

if not api_key:
    st.error("API key is not set. Please set the 'OPENAI_API_KEY' environment variable.")
    st.stop()

# Define the function to extract bibliography information using regex
def extract_bibliography_info(text):
    # Initialize default values
    bibliography_info = {
        'title': '',
        'author': '',
        'year': ''
    }

    # Clean and preprocess text
    text = text.strip().replace('\n', ' ').replace('  ', ' ')

    # Regex patterns to extract bibliographic information
    author_pattern = re.compile(r'by\s*([A-Z][a-zA-Z\s,]+)', re.IGNORECASE)
    year_pattern = re.compile(r'\b(\d{4})\b')
    title_pattern = re.compile(r'“([^”]+)”|“([^”]+)”')

    # Extract author
    author_match = author_pattern.search(text)
    if author_match:
        bibliography_info['author'] = author_match.group(1).strip()

    # Extract year
    year_match = year_pattern.search(text)
    if year_match:
        bibliography_info['year'] = year_match.group(0).strip()

    # Extract title
    title_match = title_pattern.search(text)
    if title_match:
        bibliography_info['title'] = title_match.group(1).strip() or title_match.group(2).strip()

    return bibliography_info

# Function to summarize text using Hugging Face
def summarize_text(text):
    try:
        summary = summarizer(text, max_length=150, min_length=30, do_sample=False)
        return summary[0]['summary_text']
    except Exception as e:
        st.error(f"Error summarizing text: {e}")
        return ""

# Function to extract text from an image
def extract_text_from_image(image_bytes):
    try:
        image = Image.open(io.BytesIO(image_bytes))
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        st.error(f"Error extracting text from image: {e}")
        return ""

# Function to extract text from a PDF using PyMuPDF
def extract_text_from_pdf(pdf_bytes):
    try:
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
        text = ""
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            text += page.get_text("text")
        return text
    except Exception as e:
        st.error(f"Error extracting text from PDF: {e}")
        return ""

# Function to format bibliography information for download
def format_bibliography_info(bibliography_info):
    return f"Title: {bibliography_info['title']}\n\nAuthor: {bibliography_info['author']}\n\nYear: {bibliography_info['year']}"

# Function to download as Word
def download_word(content):
    doc = docx.Document()
    doc.add_paragraph(content)
    output_file = io.BytesIO()
    doc.save(output_file)
    output_file.seek(0)
    return output_file

# Function to download as PDF
def download_pdf(content):
    output_file = io.BytesIO()
    c = canvas.Canvas(output_file, pagesize=letter)
    text_object = c.beginText(40, 750)
    text_object.setFont("Helvetica", 12)
    lines = content.splitlines()
    for line in lines:
        text_object.textLine(line)
    c.drawText(text_object)
    c.showPage()
    c.save()
    output_file.seek(0)
    return output_file

# Function to download as PowerPoint
def download_ppt(content):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    textbox = slide.shapes.add_textbox(left=Inches(1), top=Inches(1), width=Inches(8), height=Inches(5))
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = content
    output_file = io.BytesIO()
    prs.save(output_file)
    output_file.seek(0)
    return output_file

# Function to download as Excel
def download_excel(content):
    df = pd.DataFrame({"Content": [content]})
    output_file = io.BytesIO()
    df.to_excel(output_file, index=False)
    output_file.seek(0)
    return output_file

# Function to convert text to image
def text_to_image(text, image_format):
    try:
        image = Image.new('RGB', (800, 600), color='white')
        draw = ImageDraw.Draw(image)
        font = ImageFont.load_default()
        draw.text((10, 10), text, fill='black', font=font)
        output_file = io.BytesIO()
        image.save(output_file, format=image_format)
        output_file.seek(0)
        return output_file
    except Exception as e:
        st.error(f"Error converting text to image: {e}")
        return None

# Initialize session state variables
if 'page' not in st.session_state:
    st.session_state['page'] = 'main'

if 'file_name' not in st.session_state:
    st.session_state['file_name'] = 'bibliography_content'

if 'theme' not in st.session_state:
    st.session_state['theme'] = 'light'

# Apply custom theme based on session state
theme_styles = {
    'light': """
    <style>
    .main-container {
        background-color: #f5f5f5;
        padding: 20px;
        border-radius: 10px;
    }
    .btn {
        background-color: #4CAF50;
        color: white;
        padding: 10px 20px;
        border: none;
        border-radius: 5px;
        cursor: pointer;
        font-size: 16px;
    }
    .btn:hover {
        background-color: #45a049;
    }
    .header {
        font-size: 36px;
        font-weight: 600;
        color: #333;
    }
    .subheader {
        font-size: 24px;
        font-weight: 500;
        color: #555;
    }
    .text-box {
        border: 2px solid #ddd;
        padding: 10px;
        border-radius: 5px;
        background-color: #fff;
        color: black;
        max-height: 400px;
        overflow-y: auto;
        white-space: pre-wrap;
        font-family: monospace;
    }
    .file-name-container {
        display: flex;
        align-items: center;
        margin-bottom: 20px;
    }
    .file-name-container input {
        flex: 1;
        padding: 10px;
        border-radius: 5px;
        border: 1px solid #ccc;
    }
    .file-name-container button {
        margin-left: 10px;
    }
    </style>
    """,
    'dark': """
    <style>
    .main-container {
        background-color: #333;
        padding: 20px;
        border-radius: 10px;
        color: #f5f5f5;
    }
    .btn {
        background-color: #4CAF50;
        color: white;
        padding: 10px 20px;
        border: none;
        border-radius: 5px;
        cursor: pointer;
        font-size: 16px;
    }
    .btn:hover {
        background-color: #45a049;
    }
    .header {
        font-size: 36px;
        font-weight: 600;
        color: #f5f5f5;
    }
    .subheader {
        font-size: 24px;
        font-weight: 500;
        color: #ddd;
    }
    .text-box {
        border: 2px solid #555;
        padding: 10px;
        border-radius: 5px;
        background-color: #444;
        color: #f5f5f5;
        max-height: 400px;
        overflow-y: auto;
        white-space: pre-wrap;
        font-family: monospace;
    }
    .file-name-container {
        display: flex;
        align-items: center;
        margin-bottom: 20px;
    }
    .file-name-container input {
        flex: 1;
        padding: 10px;
        border-radius: 5px;
        border: 1px solid #777;
    }
    .file-name-container button {
        margin-left: 10px;
    }
    </style>
    """
}

# Apply the selected theme style
st.markdown(theme_styles[st.session_state['theme']], unsafe_allow_html=True)

# Main and About page content
def main_page():
    st.title("Bibliography Extraction and File Conversion Tool")
    
    menu = ["Extract Bibliography", "Summarize Text"]
    choice = st.sidebar.selectbox("Choose an action", menu)
    
    if choice == "Extract Bibliography":
        st.subheader("Upload a file to extract bibliography information")

        uploaded_file = st.file_uploader("Choose a file", type=["pdf", "png", "jpg", "jpeg"])

        if uploaded_file is not None:
            if uploaded_file.type == "application/pdf":
                text = extract_text_from_pdf(uploaded_file.read())
            elif uploaded_file.type in ["image/png", "image/jpeg", "image/jpg"]:
                text = extract_text_from_image(uploaded_file.read())
            else:
                st.error("Unsupported file type.")
                return

            if text:
                bibliography_info = extract_bibliography_info(text)
                formatted_bibliography = format_bibliography_info(bibliography_info)

                st.markdown("### Extracted Bibliography Information")
                st.text_area("Extracted Text", value=text, height=300)

                st.markdown("### Extracted Bibliography Details")
                st.write(f"**Title:** {bibliography_info['title']}")
                st.write(f"**Author:** {bibliography_info['author']}")
                st.write(f"**Year:** {bibliography_info['year']}")

                # File name input and save button
                col1, col2 = st.columns([4, 1])

                with col1:
                    file_name_input = st.text_input("Enter a file name for download:", value=st.session_state.file_name, key="file_name_input")

                with col2:
                    if st.button("Save"):
                        st.session_state.file_name = file_name_input  # Update the session state with the entered file name
                        st.success(f"File name saved as: {st.session_state.file_name}")

                st.markdown("### Download Extracted Bibliography")
                download_option = st.selectbox("Choose a format", ["Word", "PDF", "PowerPoint", "Excel", "PNG Image", "JPG Image"])

                if st.button("Download"):
                    file_name = f"{st.session_state.file_name}"

                    if download_option == "Word":
                        output = download_word(formatted_bibliography)
                        st.download_button(label="Download as Word", data=output, file_name=f"{file_name}.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

                    elif download_option == "PDF":
                        output = download_pdf(formatted_bibliography)
                        st.download_button(label="Download as PDF", data=output, file_name=f"{file_name}.pdf", mime="application/pdf")

                    elif download_option == "PowerPoint":
                        output = download_ppt(formatted_bibliography)
                        st.download_button(label="Download as PowerPoint", data=output, file_name=f"{file_name}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

                    elif download_option == "Excel":
                        output = download_excel(formatted_bibliography)
                        st.download_button(label="Download as Excel", data=output, file_name=f"{file_name}.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

                    elif download_option == "PNG Image":
                        output = text_to_image(formatted_bibliography, "PNG")
                        if output:
                            st.download_button(label="Download as PNG Image", data=output, file_name=f"{file_name}.png", mime="image/png")

                    elif download_option == "JPG Image":
                        output = text_to_image(formatted_bibliography, "JPEG")
                        if output:
                            st.download_button(label="Download as JPG Image", data=output, file_name=f"{file_name}.jpg", mime="image/jpeg")

    elif choice == "Summarize Text":
        st.subheader("Upload a file to summarize the content")

        uploaded_file = st.file_uploader("Choose a file", type=["pdf", "png", "jpg", "jpeg"])

        if uploaded_file is not None:
            if uploaded_file.type == "application/pdf":
                text = extract_text_from_pdf(uploaded_file.read())
            elif uploaded_file.type in ["image/png", "image/jpeg", "image/jpg"]:
                text = extract_text_from_image(uploaded_file.read())
            else:
                st.error("Unsupported file type.")
                return

            if text:
                summary = summarize_text(text)

                st.markdown("### Original Text")
                st.text_area("Original Text", value=text, height=300)

                st.markdown("### Summary")
                st.text_area("Summary", value=summary, height=150)

                # File name input and save button
                col1, col2 = st.columns([4, 1])

                with col1:
                    file_name_input = st.text_input("Enter a file name for download:", value=st.session_state.file_name, key="file_name_input")

                with col2:
                    if st.button("Save"):
                        st.session_state.file_name = file_name_input  # Update the session state with the entered file name
                        st.success(f"File name saved as: {st.session_state.file_name}")

                st.markdown("### Download Summary")
                download_option = st.selectbox("Choose a format", ["Word", "PDF", "PowerPoint", "Excel", "PNG Image", "JPG Image"])

                if st.button("Download"):
                    file_name = f"{st.session_state.file_name}"

                    if download_option == "Word":
                        output = download_word(summary)
                        st.download_button(label="Download as Word", data=output, file_name=f"{file_name}.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

                    elif download_option == "PDF":
                        output = download_pdf(summary)
                        st.download_button(label="Download as PDF", data=output, file_name=f"{file_name}.pdf", mime="application/pdf")

                    elif download_option == "PowerPoint":
                        output = download_ppt(summary)
                        st.download_button(label="Download as PowerPoint", data=output, file_name=f"{file_name}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

                    elif download_option == "Excel":
                        output = download_excel(summary)
                        st.download_button(label="Download as Excel", data=output, file_name=f"{file_name}.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

                    elif download_option == "PNG Image":
                        output = text_to_image(summary, "PNG")
                        if output:
                            st.download_button(label="Download as PNG Image", data=output, file_name=f"{file_name}.png", mime="image/png")

                    elif download_option == "JPG Image":
                        output = text_to_image(summary, "JPEG")
                        if output:
                            st.download_button(label="Download as JPG Image", data=output, file_name=f"{file_name}.jpg", mime="image/jpeg")

def about_page():
    st.title("About This App")
    st.markdown("""
    ## Overview
    This application allows you to extract bibliography information and summarize text from various file formats including PDFs and images.

    ## Features
    - **Extract Bibliography**: Upload a PDF or image to extract bibliography details such as title, author, and year.
    - **Summarize Text**: Upload a file to get a summarized version of the text content.

    ## Technologies Used
    - **Streamlit**: For building the web application interface.
    - **PyMuPDF**: For extracting text from PDFs.
    - **Tesseract OCR**: For extracting text from images.
    - **Transformers**: For text summarization using Hugging Face.
    - **spaCy**: For advanced text processing.

    ## Contact
    For any questions or feedback, please reach out to [your.email@example.com](mailto:your.email@example.com).

    ## License
    This project is licensed under the MIT License. See the [LICENSE](LICENSE) file for details.
    """)

def main():
    st.sidebar.title("Navigation")
    pages = ["Main Page", "About"]
    choice = st.sidebar.radio("Select Page", pages)

    if choice == "Main Page":
        main_page()
    elif choice == "About":
        about_page()

if __name__ == "__main__":
    main()
